import os
import json
import yaml
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
import anthropic
import argparse

# Ladda API-nyckel från .env
load_dotenv(Path(__file__).parent / ".env")

# Ladda konfiguration
def load_config():
    config_path = Path(__file__).parent / "config.yaml"
    with open(config_path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)

# Ladda logg över redan processade filer
def load_log(log_path):
    if Path(log_path).exists():
        with open(log_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

# Spara logg
def save_log(log_path, log_data):
    Path(log_path).parent.mkdir(parents=True, exist_ok=True)
    with open(log_path, "w", encoding="utf-8") as f:
        json.dump(log_data, f, ensure_ascii=False, indent=2)

# Läs textinnehåll från fil
def read_file(filepath):
    suffix = Path(filepath).suffix.lower()
    try:
        if suffix == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        elif suffix == ".docx":
            from docx import Document
            doc = Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])
        elif suffix in (".ppt", ".pptx"):
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        elif suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(filepath)
            return "\n".join([page.extract_text() or "" for page in reader.pages])
        else:
            return None
    except Exception as e:
        print(f"  Kunde inte läsa {filepath}: {e}")
        return None

# Kontrollera om filen är låst (öppen i annat program)
def check_file_locked(filepath):
    try:
        with open(filepath, "a"):
            pass
        return False
    except IOError:
        return True

# Analysera dokument med Claude
def analyze_document(client, model, max_tokens, filepath, content, default_author=""):
    prompt = f"""Analysera följande dokument och svara ENDAST med ett JSON-objekt i detta exakta format (inga kommentarer, inga markdown-kodblock):
{{
  "title": "dokumentets titel eller ett beskrivande namn om titel saknas",
  "author": "författare i formatet 'Efternamn, Förnamn' eller 'Okänd' om det inte framgår. Flera författare separeras med semikolon",
  "summary": "sammanfattning på svenska med 30-150 ord beroende på innehållets komplexitet",
  "type": "en av: artikel, uppsats, bok, predikan, studie, övrigt (använd 'bok' även för äldre böcker utan ISBN)",
  "year": "utgivningsår som heltal eller null om okänt",
  "date_full": "exakt datum i formatet YYYY-MM-DD om det går att fastställa, annars null",
  "is_citable": true eller false,

  "publication": "tidskrift eller bok som artikeln publicerats i, eller null om ej artikel",

  "publisher": "förlagets namn eller null om okänt eller ej bok",
  "publisher_place": "utgivningsort eller null om okänd eller ej bok",
  "isbn": "ISBN om det finns angivet, annars null",
  "pages_total": "totalt antal sidor som heltal eller null om okänt",
  "edition": "upplaga t.ex. '2nd edition' eller null om ej angiven",

  "institution": "lärosäte för uppsats/avhandling eller null om ej tillämpligt",
  "institution_place": "ort för lärosätet eller null om ej tillämpligt",
  "thesis_type": "t.ex. 'Kandidatuppsats', 'Masteruppsats', 'Doktorsavhandling' eller null om ej tillämpligt"
}}

Filnamn: {Path(filepath).name}

Dokumentets innehåll (kan vara avkortat):
{content[:6000]}
"""
    message = client.messages.create(
        model=model,
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}]
    )
    raw = message.content[0].text.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    result = json.loads(raw.strip())
    if default_author and result.get("author", "").lower() in ("okänd", "unknown", ""):
        result["author"] = default_author
    return result

# Hitta alla filer att processa
def find_files(folders, extensions, log):
    files = []
    for folder in folders:
        for root, dirs, filenames in os.walk(folder):
            # Hoppa över analyzer-mappar
            dirs[:] = [d for d in dirs if d != "analyzer"]
            for filename in filenames:
                filepath = str(Path(root) / filename)
                if Path(filepath).suffix.lower() in extensions:
                    if filepath not in log:
                        files.append(filepath)
                    else:
                        print(f"  Hoppar över (redan processad): {filename}")
    return files

# Generera Word-rapport
def generate_word_report(results, output_path, folder_name=""):
    from docx import Document as DocxDocument
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    if check_file_locked(output_path):
        print(f"\n⚠️  Kan inte spara rapporten – filen är öppen i Word:")
        print(f"   {output_path}")
        print(f"   Stäng filen och tryck Enter för att försöka igen...")
        input()
    doc = DocxDocument()

    # Titel
    title = doc.add_heading(f"Analys av innehåll i mappen {folder_name}", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Datum
    date_para = doc.add_paragraph(f"Genererad: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    # Sammanfattning
    doc.add_heading(f"Totalt analyserade dokument: {len(results)}", level=2)
    doc.add_paragraph()

    # Gruppera efter typ
    by_type = {}
    for r in results:
        t = r.get("type", "övrigt")
        by_type.setdefault(t, []).append(r)

    for doc_type, items in sorted(by_type.items()):
        doc.add_heading(doc_type.capitalize(), level=1)
        for item in items:
            # Rubrik
            doc.add_heading(item.get("title", "Utan titel"), level=2)

            # Författare och år/datum
            date_str = item.get("date_full") or item.get("year") or "okänt"
            meta = f"Författare: {item.get('author', 'Okänd')}  |  År: {date_str}"
            p = doc.add_paragraph(meta)
            p.runs[0].italic = True

            # Artikelspecifikt: publikation
            if item.get("type") == "artikel" and item.get("publication"):
                pub = doc.add_paragraph(f"Publikation: {item['publication']}")
                pub.runs[0].italic = True

            # Uppsatsspecifikt: lärosäte
            if item.get("type") in ("uppsats", "avhandling") and item.get("institution"):
                inst_str = item["institution"]
                if item.get("institution_place"):
                    inst_str += f", {item['institution_place']}"
                if item.get("thesis_type"):
                    inst_str += f" ({item['thesis_type']})"
                ins = doc.add_paragraph(f"Lärosäte: {inst_str}")
                ins.runs[0].italic = True

            # Filnamn
            fn = doc.add_paragraph(f"Fil: {Path(item['filepath']).name}")
            fn.runs[0].font.size = Pt(9)
            fn.runs[0].font.color.rgb = RGBColor(128, 128, 128)

            # Undermapp (om filen inte ligger direkt i rotmappen)
            file_folder = Path(item['filepath']).parent.name
            if file_folder != folder_name:
                mapp = doc.add_paragraph(f"Mapp: {file_folder}")
                mapp.runs[0].font.size = Pt(9)
                mapp.runs[0].font.color.rgb = RGBColor(128, 128, 128)

            # Sammanfattning
            doc.add_paragraph(item.get("summary", ""))
            doc.add_paragraph()

    doc.save(output_path)
    print(f"Word-rapport sparad: {output_path}")

# Formatera författarnamn för RIS (efternamn, förnamn)
def format_ris_author(name):
    name = name.strip()
    if "," in name:
        return name
    parts = name.split()
    if len(parts) >= 2:
        return f"{parts[-1]}, {' '.join(parts[:-1])}"
    return name

# Omvandla filväg till file:/// URI
# Används för att skapa länkar till bilagor i Zotero-rapporten
def filepath_to_uri(filepath):
    # Omvandla Windows-sökväg till file:/// URI
    path = Path(filepath).resolve()
    # Ersätt bakåtsnedstreck med framåtsnedstreck
    uri = path.as_uri()
    return uri

# Generera Zotero RIS-export
def generate_zotero_export(results, output_path):
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    NON_CITABLE_TYPES = {"predikan", "övrigt"}
    NON_CITABLE_EXTENSIONS = {".ppt", ".pptx"}
    
    citable = [
        r for r in results
        if r.get("is_citable")
        and r.get("type") not in NON_CITABLE_TYPES
        and Path(r.get("filepath", "")).suffix.lower() not in NON_CITABLE_EXTENSIONS
    ]

    type_map = {
        "artikel": "JOUR",
        "uppsats": "THES",
        "avhandling": "THES",
        "bok": "BOOK",
        "studie": "RPRT",
    }

    lines = []
    for item in citable:
        ris_type = type_map.get(item.get("type", ""), "GEN")
        lines.append(f"TY  - {ris_type}")
        lines.append(f"TI  - {item.get('title', 'Utan titel')}")
        if item.get("filepath"):
            lines.append(f"L1  - {filepath_to_uri(item['filepath'])}")

        author = item.get("author", "Okänd")
        for a in author.split(";"):
            lines.append(f"AU  - {format_ris_author(a)}")

        # Datum
        if item.get("date_full"):
            lines.append(f"DA  - {item['date_full']}")
        elif item.get("year"):
            lines.append(f"PY  - {item['year']}")

        # Sammanfattning
        if item.get("summary"):
            lines.append(f"N2  - {item['summary']}")

        # Artikelspecifikt
        if item.get("publication"):
            lines.append(f"JO  - {item['publication']}")

        # Bokspecifikt
        if item.get("publisher"):
            lines.append(f"PB  - {item['publisher']}")
        if item.get("publisher_place"):
            lines.append(f"CY  - {item['publisher_place']}")
        if item.get("isbn"):
            lines.append(f"SN  - {item['isbn']}")
        if item.get("pages_total"):
            lines.append(f"SP  - {item['pages_total']} sidor")
        if item.get("edition"):
            lines.append(f"ET  - {item['edition']}")

        # Uppsatsspecifikt
        if item.get("institution"):
            lines.append(f"PB  - {item['institution']}")
        if item.get("institution_place"):
            lines.append(f"CY  - {item['institution_place']}")
        if item.get("thesis_type"):
            lines.append(f"M3  - {item['thesis_type']}")

        lines.append("ER  - ")
        lines.append("")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    print(f"Zotero RIS-fil sparad: {output_path} ({len(citable)} poster)")


# Huvudfunktion
def main():
    # Hantera kommandoradsargument
    parser = argparse.ArgumentParser(description="Analysera dokument i en mapp")
    parser.add_argument("--folder", type=str, help="Mapp att analysera (överskriver config.yaml)")
    parser.add_argument("--noris", action="store_true", help="Skapa ingen Zotero RIS-fil")
    parser.add_argument("--refresh", action="store_true", help="Radera loggen och analysera allt från scratch")
    args = parser.parse_args()

    config = load_config()

    # Använd default_author från config
    default_author = config.get("default_author", "Okänd")

    # Överskrid config om --folder angivits
    if args.folder:
        config["folders"] = [str(Path(args.folder).resolve())]

    # Definiera alla sökvägar tidigt
    folder_name = Path(config["folders"][0]).name
    base_output = Path(config["folders"][0]) / "analyzer"
    log_path = str(base_output / "processed_files.json")
    if args.refresh:
        if Path(log_path).exists():
            Path(log_path).unlink()
            print("Logg raderad - analyserar allt från scratch.")
        log = {}
    else:
        log = load_log(log_path)
    report_path = str(base_output / f"analys-{folder_name}.docx")
    zotero_path = str(base_output / f"zotero_import_{folder_name}.ris")

    log = load_log(log_path)

    client = anthropic.Anthropic()
    model = config["anthropic"]["model"]
    max_tokens = config["anthropic"]["max_tokens"]

    print(f"Startar analys: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    print(f"Redan processade filer: {len(log)}")

    files = find_files(config["folders"], config["extensions"], log)
    print(f"Nya filer att processa: {len(files)}\n")

    results = []

    for i, filepath in enumerate(files, 1):
        filename = Path(filepath).name
        print(f"[{i}/{len(files)}] Analyserar: {filename}")

        content = read_file(filepath)
        if not content or len(content.strip()) < 50:
            print(f"  Hoppar över – tomt eller oläsbart innehåll")
            continue

        try:
            analysis = analyze_document(client, model, max_tokens, filepath, content, default_author)
            analysis["filepath"] = filepath
            results.append(analysis)

            log[filepath] = {
                "processed": datetime.now().isoformat(),
                "title": analysis.get("title"),
                "author": analysis.get("author"),
                "analysis": analysis
            }
            save_log(log_path, log)
            print(f"  ✓ {analysis.get('author', 'Okänd')} – {analysis.get('title', 'Utan titel')}")

        except Exception as e:
            print(f"  ✗ Fel vid analys: {e}")

    print(f"\nKlart! {len(results)} dokument analyserade.")

    # Bygg lista med alla resultat - nya + tidigare analyserade
    all_results = []
    for filepath, entry in log.items():
        if "analysis" in entry:
            all_results.append(entry["analysis"])

    print(f"Totalt i rapport: {len(all_results)} dokument.")

    generate_word_report(all_results, report_path, folder_name)
    generate_word_report(all_results, report_path, folder_name)
    if not args.noris:
        generate_zotero_export(all_results, zotero_path)

if __name__ == "__main__":
    main()